import streamlit as st  
from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE  
from PIL import Image  
from io import BytesIO  
import requests  
import base64  
from docx import Document  
from docx.shared import Pt  
import fitz  # PyMuPDF  
import os  
import cv2  
import numpy as np  
import io  
  
# Azure OpenAI credentials  
azure_endpoint = "https://gpt-4omniwithimages.openai.azure.com/"  
api_key = "6e98566acaf24997baa39039b6e6d183"  
api_version = "2024-02-01"  
model = "GPT-40-mini"  
  
# URL of your Azure function endpoint  
azure_function_url = 'https://doc2pdf.azurewebsites.net/api/HttpTrigger1'  
  
# Function to convert PPT to PDF using Azure Function  
def ppt_to_pdf(ppt_file, pdf_file):  
    mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'  
    headers = {  
        "Content-Type": "application/octet-stream",  
        "Content-Type-Actual": mime_type  
    }  
    with open(ppt_file, 'rb') as file:  
        response = requests.post(azure_function_url, data=file.read(), headers=headers)  
        if response.status_code == 200:  
            with open(pdf_file, 'wb') as pdf_out:  
                pdf_out.write(response.content)  
            return True  
        else:  
            st.error(f"File conversion failed with status code: {response.status_code}")  
            st.error(f"Response: {response.text}")  
            return False  
  
# Function to encode image as base64  
def encode_image(image):  
    return base64.b64encode(image).decode("utf-8")  
  
def get_image_explanation(base64_image):  
    headers = {  
        "Content-Type": "application/json",  
        "api-key": api_key  
    }  
    data = {  
        "model": model,  
        "messages": [  
            {"role": "system", "content": "You are a helpful assistant that responds in Markdown."},  
            {"role": "user", "content": [  
                {"type": "text", "text": "Explain the content of this image in a single, coherent paragraph. The explanation should be concise and semantically meaningful, summarizing all major points from the image in one paragraph. Avoid using bullet points or separate lists."},  
                {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}  
            ]}  
        ],  
        "temperature": 0.7  
    }  
  
    response = requests.post(  
        f"{azure_endpoint}/openai/deployments/{model}/chat/completions?api-version={api_version}",  
        headers=headers,  
        json=data  
    )  
  
    if response.status_code == 200:  
        result = response.json()  
        return result["choices"][0]["message"]["content"]  
    else:  
        st.error(f"Error: {response.status_code} - {response.text}")  
        return None  
  
def extract_text_from_ppt(ppt_file):  
    presentation = Presentation(ppt_file)  
    text_content = []  
    for slide_number, slide in enumerate(presentation.slides, start=1):  
        slide_text = []  
        for shape in slide.shapes:  
            if shape.has_text_frame:  
                for paragraph in shape.text_frame.paragraphs:  
                    for run in paragraph.runs:  
                        slide_text.append(run.text)  
        slide_title = slide.shapes.title.text if slide.shapes.title else "Untitled Slide"  
        text_content.append({"slide_number": slide_number, "slide_title": slide_title, "text": " ".join(slide_text)})  
    return text_content  
  
def is_image_of_interest(shape):  
    """Check if a shape contains an image in formats of interest"""  
    try:  
        if hasattr(shape, "image"):  
            image_ext = os.path.splitext(shape.image.filename)[1].lower()  
            if image_ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff"]:  
                return image_ext  
    except Exception:  
        pass  
    return None  
  
def detect_image_slides(ppt_bytes):  
    """Detect slides containing images in the desired formats"""  
    ppt = Presentation(io.BytesIO(ppt_bytes))  
    image_slides = {}  
    for i, slide in enumerate(ppt.slides):  
        for shape in slide.shapes:  
            image_format = is_image_of_interest(shape)  
            if image_format:  
                slide_number = i + 1  
                image_slides[slide_number] = image_format  
                break  
    return image_slides  
  
def identify_visual_elements(ppt_bytes):  
    """Identify slides with visual elements"""  
    presentation = Presentation(io.BytesIO(ppt_bytes))  
    visual_slides = []  
    for slide_number, slide in enumerate(presentation.slides, start=1):  
        has_visual_elements = False  
        for shape in slide.shapes:  
            if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.TABLE, MSO_SHAPE_TYPE.CHART,  
                                    MSO_SHAPE_TYPE.GROUP, MSO_SHAPE_TYPE.AUTO_SHAPE}:  
                has_visual_elements = True  
                break  
        if has_visual_elements:  
            visual_slides.append(slide_number)  
    return visual_slides  
  
def combine_slide_numbers(image_slides, visual_slides):  
    """Combine slide numbers from image slides and visual element slides"""  
    combined_slides = set(image_slides.keys()).union(set(visual_slides))  
    return sorted(list(combined_slides))  
  
def capture_slide_images(pdf_file, slide_numbers):  
    """Capture images from identified slides in the PDF"""  
    doc = fitz.open(pdf_file)  
    images = []  
    for slide_number in slide_numbers:  
        page = doc[slide_number - 1]  
        pix = page.get_pixmap()  
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)  
        buffer = BytesIO()  
        img.save(buffer, format="PNG")  
        images.append({"slide_number": slide_number, "image": buffer.getvalue()})  
    return images  
  
def generate_text_insights(text_content, visual_slides, text_length):  
    headers = {  
        "Content-Type": "application/json",  
        "api-key": api_key  
    }  
    insights = []  
  
    # Set temperature based on text_length  
    if text_length == "Standard":  
        temperature = 0.3  
    elif text_length == "Blend":  
        temperature = 0.5  
    elif text_length == "Creative":  
        temperature = 0.7  
  
    for slide in text_content:  
        slide_text = slide['text']  
        slide_number = slide['slide_number']  
        prompt = f"""  
        NOTE: If a slide contains less than 30 words and adds no meaningful value, skip generating content and state, "This slide has no quality content." Ensure accurate validation by comparing the slide's content to others, and only dismiss content that lacks quality or relevance.  
        I want you to immediately begin with one of the following phrases based on the slide title:  
        (a) If and only if the slide title contains the keyword "Background," begin the explanation with "The prior solutions include..." Proceed by discussing only the prior solution presented in the slide. Ensure no mention of any proposal or disclosure occurs at this stage, and strictly limit the explanation to the prior solutions.  
        (b) If and only if the slide title contains the keyword "Proposal," start the explanation with "The present disclosure includes..." Focus exclusively on discussing the proposal or invention presented in the slide. Ensure that no background information is referenced, and strictly adhere to the proposal/invention-related content.  
        (c) If the slide title does not contain either "Background" or "Proposal," start the explanation with "Aspects of the present disclosure include..." Discuss the key aspects of the slide's content, ensuring no mention of prior solutions or proposals. Adhere to the neutral tone, focusing on the core aspects of the slide's content. The information should be delivered directly and engagingly in a single, coherent paragraph. Avoid phrases like 'The slide presents,' 'discusses,' 'outlines,' or 'content.' The explanation should be concise and semantically meaningful, summarizing all major points in one paragraph without bullet points. The text should adhere to the following style guidelines:  
        1. Remove all listed profanity words.  
        2. Use passive voice.  
        3. Use conditional and tentative language, such as "may include," "in some aspects," and "aspects of the present disclosure."  
        4. Replace "Million" with "1,000,000" and "Billion" with "1,000,000,000."  
        5. Maintain the following tone characteristics: Precision and Specificity, Formality, Complexity, Objective and Impersonal, Structured and Systematic.  
        6. Follow these style elements: Formal and Objective, Structured and Systematic, Technical Jargon and Terminology, Detailed and Specific, Impersonal Tone, Instructional and Descriptive, Use of Figures and Flowcharts, Legal and Protective Language, Repetitive and Redundant, Examples and Clauses.  
        7. Use the following conditional and tentative language phrases: may include, in some aspects, aspects of the present disclosure, by way of example, may be, may further include, may be used, may occur, may use, may monitor, may periodically wake up, may demodulate, may consume, can be performed, may enter and remain, may correspond to, may also include, may be identified in response to, may be further a function of, may be multiplied by, may schedule, may select, may also double, may further comprise, may be configured to, may correspond to a duration value, may correspond to a product of, may be closer, may be significant, may not be able, may result, may reduce, may be operating in, may further be configured to, may further process, may be executed by, may be received, may avoid, may indicate, may be selected, may be proactive, may perform, may be necessary, may be amplified, may involve, may require, may be stored, may be accessed, may be transferred, may be implemented, may include instructions to, may depend upon, may communicate, may be generated, may be configured.  
        8. Maintain the exact wording in the generated content. Do not substitute words with synonyms. For example, "instead" should remain "instead" and not be replaced with "conversely."  
        9. Replace the phrase "further development" with "our disclosure" in all generated content.  
        10. Make sure to use LaTeX formatting for all mathematical symbols, equations, subscripting, and superscripting to ensure they are displayed correctly in the output.  
        11. When encountering programmatic terms or equations, ensure they are accurately represented and contextually retained.  
        {slide_text}  
        """  
        if text_length == "Standard":  
            prompt += "\n\nGenerate a short paragraph."  
        elif text_length == "Blend":  
            prompt += "\n\nGenerate a medium-length paragraph."  
        elif text_length == "Creative":  
            prompt += "\n\nGenerate a longer paragraph."  
  
        data = {  
            "model": model,  
            "messages": [{"role": "system", "content": "You are a helpful assistant."}, {"role": "user", "content": prompt}],  
            "temperature": temperature  
        }  
  
        response = requests.post(  
            f"{azure_endpoint}/openai/deployments/{model}/chat/completions?api-version={api_version}",  
            headers=headers,  
            json=data  
        )  
  
        if response.status_code == 200:  
            result = response.json()  
            insights.append({"slide_number": slide['slide_number'], "slide_title": slide['slide_title'], "insight": result["choices"][0]["message"]["content"]})  
        else:  
            st.error(f"Error: {response.status_code} - {response.text}")  
            insights.append({"slide_number": slide['slide_number'], "slide_title": slide['slide_title'], "insight": "Error in generating insight"})  
  
    return insights  
  
def generate_image_insights(image_content, text_length, api_key, azure_endpoint, model, api_version):  
    insights = []  
    low_quality_slides = []  
  
    # Set temperature based on text_length  
    if text_length == "Standard":  
        temperature = 0.3  
    elif text_length == "Blend":  
        temperature = 0.5  
    elif text_length == "Creative":  
        temperature = 0.7  
  
    for image_data in image_content:  
        base64_image = encode_image(image_data['image'])  
        slide_number = image_data['slide_number']  
  
        # Determine how many images are on the slide  
        images_on_slide = [img for img in image_content if img['slide_number'] == slide_number]  
        image_ref = f"figure {slide_number}"  
        if len(images_on_slide) > 1:  
            image_ref += f"({chr(97 + images_on_slide.index(image_data))})"  
  
        headers = {  
            "Content-Type": "application/json",  
            "api-key": api_key  
        }  
  
        prompt = f"""  1. Step 1: Detecting and Listing Visual Elements  
   Begin by identifying all visual elements (e.g., figures, images, charts, graphs, and tables) within the document. Ensure that no visual element is overlooked, especially when multiple elements are arranged together.  
   - Step 1(a): If there are multiple visual elements, list them using the format: "Referring to {image_ref}(1), {image_ref}(2), etc." Ensure that each element is referenced correctly.  
   - Step 1(b): After listing, provide a brief description for each visual element based on its position and context within the document, ensuring that no element is skipped.

2. Step 2: Structuring the Explanation  
   After identifying the visual elements, structure the explanation based on the document's content and the type of information presented. Use appropriate phrasing for different sections:
   - If the content discusses background information or existing solutions, begin with: "The prior solutions include..." and limit the explanation to only the background or existing knowledge.
   - If the document introduces a new proposal or solution, start with: "The present disclosure includes..." and focus only on the proposed solution or invention.
   - If neither background nor a proposal is being discussed, begin with: "Aspects of the present content include..." and focus on key points from the document, avoiding any mention of prior solutions or proposals.

3. Step 3: Explaining Examples and Keywords  
   For any instance where the document includes an example or refers to a specific case (e.g., using the word "example" or "e.g."), ensure that:
   - The term "example" is always reproduced exactly as it appears in the document.
   - Every example is explained thoroughly.
   - If there are multiple examples, ensure each one is detailed individually without merging them into a single sentence.

4. Step 4: Graphs and Charts  
   For any visual elements like graphs or charts, provide an explanation that includes:
   - A detailed description of the axes (x-axis and y-axis) and any labeled data points.
   - The overall meaning and trends indicated by the graph.

5. Step 5: Perspective Views  
   For images or figures featuring a perspective view, start by identifying that the image is a perspective view. Then, provide a detailed explanation of the spatial relationships, angles, and depth depicted in the image.

6. Step 6: Content Integration  
   Integrate textual content with visual elements:
   - Begin by analyzing the text in the document, accurately reproducing key points, and ensuring the context is preserved.
   - Smoothly integrate explanations of the images or figures into the flow of the text, ensuring that the relationship between the textual content and visual elements is clear.

7. Step 7: Style and Tone Guidelines  
   Adhere to the following style guidelines:
   - Use passive voice consistently throughout the explanation.
   - Replace large numerical values with their full numeric form (e.g., "1,000,000" instead of "Million").
   - Avoid the verb "consist" and use precise technical language (e.g., "defined as," "includes," "comprises").
   - Use conditional and tentative phrases such as "may include," "in some aspects," and "by way of example" where appropriate.
   - Maintain a formal tone, technical jargon, and an impersonal voice.

8. Step 8: Avoid Unnecessary Phrasing  
   Do not start sentences with phrases like "The document," "The text," or "The figure." Instead, ensure a natural flow in the explanation by directly addressing the content and context.

9. Step 9: Follow Content-Specific Logic  
   Follow the logic inherent in the document’s structure. For instance:
   - If the document follows a step-by-step explanation, mirror that structure in your output.
   - Ensure accurate representation and integration of any equations, flowcharts, or complex data visualizations that appear within the document.

10. Step 10: Cohesive Paragraph for Points with Examples  
    For slides or sections with bullet points and examples, generate a cohesive paragraph that integrates these points while ensuring that all examples are specifically called out and fully explained. Always use the word "example" each time it appears in the content, and maintain clarity and continuity in the overall description.
"""  
        data = {  
            "model": model,  
            "messages": [  
                {"role": "system", "content": "You are a helpful assistant that responds in Markdown."},  
                {"role": "user", "content": [  
                    {"type": "text", "text": prompt},  
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}]  
                }  
            ],  
            "temperature": temperature  
        }  
  
        response = requests.post(  
            f"{azure_endpoint}/openai/deployments/{model}/chat/completions?api-version={api_version}",  
            headers=headers,  
            json=data  
        )  
  
        if response.status_code == 200:  
            result = response.json()  
            insight = result["choices"][0]["message"]["content"]  
  
            if "This slide has no quality content." in insight:  
                low_quality_slides.append(slide_number)  
            else:  
                insights.append({  
                    "slide_number": slide_number,  
                    "slide_title": image_data.get('slide_title', 'Untitled Slide'),  
                    "insight": insight  
                })  
        else:  
            st.error(f"Error: {response.status_code} - {response.text}")  
            insights.append({  
                "slide_number": slide_number,  
                "slide_title": image_data.get('slide_title', 'Untitled Slide'),  
                "insight": "Error in generating insight"  
            })  
  
    return insights, low_quality_slides   
  
def aggregate_content(text_insights, image_insights):  
    aggregated_content = []  
    for text in text_insights:  
        slide_number = text['slide_number']  
        slide_title = text['slide_title']  
        text_insight = text['insight']  
        image_insight = next((img['insight'] for img in image_insights if img['slide_number'] == slide_number), None)  
        if image_insight:  
            aggregated_content.append({  
                "slide_number": slide_number,  
                "slide_title": slide_title,  
                "content": f"{image_insight}"  
            })  
        else:  
            aggregated_content.append({  
                "slide_number": slide_number,  
                "slide_title": slide_title,  
                "content": text_insight  
            })  
    return aggregated_content  
  
def sanitize_text(text):  
    if text:  
        sanitized = ''.join(c for c in text if c.isprintable() and c not in {'\x00', '\x01', '\x02', '\x03', '\x04', '\x05', '\x06', '\x07', '\x08', '\x0B', '\x0C', '\x0E', '\x0F', '\x10', '\x11', '\x12', '\x13', '\x14', '\x15', '\x16', '\x17', '\x18', '\x19', '\x1A', '\x1B', '\x1C', '\x1D', '\x1E', '\x1F'})  
        return sanitized  
    return text  
  
def ensure_proper_spacing(text):  
    if text:  
        return text.replace('. ', '. ').replace('. ', '. ')  
    return text  
  
def save_content_to_word(aggregated_content, output_file_name, extracted_images, low_quality_slides):  
    doc = Document()  
    style = doc.styles['Normal']  
    font = style.font  
    font.name = 'Times New Roman'  
    font.size = Pt(10.5)  # Reduced font size for paragraphs  
    paragraph_format = style.paragraph_format  
    paragraph_format.line_spacing = 1.5  
    paragraph_format.alignment = 3  # Justify  
  
    for slide in aggregated_content:  
        sanitized_title = sanitize_text(slide['slide_title'])  
        sanitized_content = sanitize_text(slide['content'])  
        properly_spaced_content = ensure_proper_spacing(sanitized_content)  
        doc.add_heading(f"[[{slide['slide_number']}, {sanitized_title}]]", level=1)  
        if properly_spaced_content:  # Only add content if it exists  
            doc.add_paragraph(properly_spaced_content)  
  
    # Add extracted images after the generated content  
    if extracted_images:  
        doc.add_heading("Extracted Images", level=1)  
        for idx, (image, slide_number) in enumerate(extracted_images):  
            if slide_number not in low_quality_slides:  # Skip low-quality slides  
                _, buffer = cv2.imencode('.png', image)  
                image_stream = BytesIO(buffer)  
                doc.add_paragraph(f"Image from Slide {slide_number}:")  
                doc.add_picture(image_stream, width=doc.sections[0].page_width - doc.sections[0].left_margin - doc.sections[0].right_margin)  
                doc.add_paragraph("\n")  # Add space after image  
  
    output = BytesIO()  
    doc.save(output)  
    output.seek(0)  
    return output  
  
def extract_and_clean_page_image(page, top_mask, bottom_mask, left_mask, right_mask):  
    # Get the page as an image  
    pix = page.get_pixmap()  
    img_data = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)  
  
    # Convert the image to BGR format for OpenCV  
    img_bgr = cv2.cvtColor(img_data, cv2.COLOR_RGB2BGR)  
  
    # Convert to grayscale for processing  
    gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)  
  
    # Threshold the image to get binary image  
    _, binary = cv2.threshold(gray, 240, 255, cv2.THRESH_BINARY_INV)  
  
    # Detect contours of possible images/diagrams  
    contours, _ = cv2.findContours(binary, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)  
  
    # Check if there are any valid contours (image regions)  
    valid_contours = [cv2.boundingRect(contour) for contour in contours if cv2.boundingRect(contour)[2] > 50 and cv2.boundingRect(contour)[3] > 50]  
    if not valid_contours:  
        return None  # Skip the page if no valid images/diagrams are found  
  
    # Create a mask for the detected contours  
    mask = np.zeros_like(binary)  
    for x, y, w, h in valid_contours:  
        # Apply the adjustable top, bottom, left, and right masking values from the sliders  
        # Ensure coordinates do not go out of image bounds  
        x1 = max(0, x - left_mask)  
        y1 = max(0, y - top_mask)  
        x2 = min(img_bgr.shape[1], x + w + right_mask)  
        y2 = min(img_bgr.shape[0], y + h + bottom_mask)  
        cv2.rectangle(mask, (x1, y1), (x2, y2), 255, -1)  
  
    # Use the mask to keep only the regions with images/diagrams  
    text_removed = cv2.bitwise_and(img_bgr, img_bgr, mask=mask)  
  
    # Set the background to white where the mask is not applied  
    background = np.ones_like(img_bgr) * 255  
    cleaned_image = np.where(mask[:, :, None] == 255, text_removed, background)  
  
    # Convert cleaned image to grayscale  
    cleaned_image_gray = cv2.cvtColor(cleaned_image, cv2.COLOR_BGR2GRAY)  
    return cleaned_image_gray  
  
def extract_images_from_pdf(pdf_file, top_mask, bottom_mask, left_mask, right_mask):  
    # Open the PDF file  
    pdf_document = fitz.open(pdf_file)  
    page_images = []  
  
    for page_num in range(len(pdf_document)):  
        page = pdf_document.load_page(page_num)  
  
        # Extract and clean the page image  
        cleaned_image = extract_and_clean_page_image(page, top_mask, bottom_mask, left_mask, right_mask)  
        if cleaned_image is not None:  
            page_images.append((cleaned_image, page_num + 1))  # Keep track of the slide number  
  
    pdf_document.close()  
    return page_images  
  
def main():  
    st.title("PPT Insights Extractor")  
  
    text_length = st.select_slider(  
        "Content Generation Slider",  
        options=["Standard", "Blend", "Creative"],  
        value="Blend"  
    )  
  
    # Add Title and Information Button for Image Extraction Slider  
    st.sidebar.markdown("### Image Extraction Slider")  
  
    # Initialize session state variables for the sliders  
    if 'top_mask' not in st.session_state:  
        st.session_state.top_mask = 40  
    if 'bottom_mask' not in st.session_state:  
        st.session_state.bottom_mask = 40  
    if 'left_mask' not in st.session_state:  
        st.session_state.left_mask = 85  
    if 'right_mask' not in st.session_state:  
        st.session_state.right_mask = 85  
  
    # Arrange the buttons in a row using columns  
    col1, col2 = st.sidebar.columns(2)  
    with col1:  
        if st.button("Default"):  
            st.session_state.top_mask = 40  
            st.session_state.bottom_mask = 40  
            st.session_state.left_mask = 85  
            st.session_state.right_mask = 85  
  
    with col2:  
        if st.button("A4"):  
            st.session_state.top_mask = 70  
            st.session_state.bottom_mask = 70  
            st.session_state.left_mask = 85  
            st.session_state.right_mask = 85  
  
    # Add sliders to adjust the top, bottom, left, and right masking values  
    top_mask = st.sidebar.slider("Adjust Top Masking Value", min_value=10, max_value=100, value=st.session_state.top_mask, step=1)  
    bottom_mask = st.sidebar.slider("Adjust Bottom Masking Value", min_value=10, max_value=100, value=st.session_state.bottom_mask, step=1)  
    left_mask = st.sidebar.slider("Adjust Left Masking Value", min_value=10, max_value=500, value=st.session_state.left_mask, step=1)  
    right_mask = st.sidebar.slider("Adjust Right Masking Value", min_value=10, max_value=200, value=st.session_state.right_mask, step=1)  
  
    # Update session state if sliders are moved  
    if top_mask != st.session_state.top_mask or bottom_mask != st.session_state.bottom_mask or left_mask != st.session_state.left_mask or right_mask != st.session_state.right_mask:  
        st.session_state.top_mask = top_mask  
        st.session_state.bottom_mask = bottom_mask  
        st.session_state.left_mask = left_mask  
        st.session_state.right_mask = right_mask  
  
    uploaded_ppt = st.file_uploader("Upload a PPT file", type=["pptx"])  
  
    if uploaded_ppt is not None:  
        st.info("Processing PPT file...")  
  
        # Extract the base name of the uploaded PPT file  
        ppt_filename = uploaded_ppt.name  
        base_filename = os.path.splitext(ppt_filename)[0]  
        output_word_filename = f"{base_filename}.docx"  
  
        try:  
            # Convert PPT to PDF  
            with open("temp_ppt.pptx", "wb") as f:  
                f.write(uploaded_ppt.read())  
            if not ppt_to_pdf("temp_ppt.pptx", "temp_pdf.pdf"):  
                st.error("PDF conversion failed. Please check the uploaded PPT file.")  
                return  
  
            # Read the PPT file content as bytes  
            with open("temp_ppt.pptx", "rb") as f:  
                ppt_bytes = f.read()  
  
            # Extract text and identify slides with visual elements  
            text_content = extract_text_from_ppt(BytesIO(ppt_bytes))  
            visual_slides = identify_visual_elements(ppt_bytes)  
  
            # Detect slides with images  
            image_slides = detect_image_slides(ppt_bytes)  
  
            # Combine slide numbers from both functions  
            combined_slides = combine_slide_numbers(image_slides, visual_slides)  
  
            # Capture images of marked slides  
            slide_images = capture_slide_images("temp_pdf.pdf", combined_slides)  
  
            st.info("Generating text insights...")  
            text_insights = generate_text_insights(text_content, visual_slides, text_length)  
  
            st.info("Generating image insights...")  
            image_insights, low_quality_slides = generate_image_insights(slide_images, text_length, api_key, azure_endpoint, model, api_version)  
  
            # Filter out low-quality slides from text and image insights  
            filtered_text_insights = [insight for insight in text_insights if insight['slide_number'] not in low_quality_slides]  
            filtered_image_insights = [insight for insight in image_insights if insight['slide_number'] not in low_quality_slides]  
  
            st.info("Extracting additional images...")  
            extracted_images = extract_images_from_pdf("temp_pdf.pdf", top_mask, bottom_mask, left_mask, right_mask)  
  
            st.info("Aggregating content...")  
            aggregated_content = aggregate_content(filtered_text_insights, filtered_image_insights)  
  
            st.info("Saving to Word document...")  
            output_doc = save_content_to_word(aggregated_content, output_word_filename, extracted_images, low_quality_slides)  
  
            st.download_button(label="Download Word Document", data=output_doc, file_name=output_word_filename)  
  
            st.success("Processing completed successfully!")  
        except Exception as e:  
            st.error(f"An error occurred: {e}")  
  
if __name__ == "__main__":  
    main()  
