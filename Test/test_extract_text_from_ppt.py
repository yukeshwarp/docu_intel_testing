import pytest  
from unittest.mock import patch  
from pptx import Presentation  
from pptx.util import Inches  
  
# Assuming the extract_text_from_ppt function is in a module named ppt_extractor  
from docu_intel import extract_text_from_ppt  
  
@pytest.fixture  
def mock_presentation():  
    prs = Presentation()  
      
    # Slide 1  
    slide_1 = prs.slides.add_slide(prs.slide_layouts[1])  
    title_1 = slide_1.shapes.title  
    title_1.text = "Slide 1 Title"  
    content_1 = slide_1.shapes.placeholders[1]  
    content_1.text = "This is the first slide."  
  
    # Slide 2  
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  
    title_2 = slide_2.shapes.title  
    title_2.text = "Slide 2 Title"  
    content_2 = slide_2.shapes.placeholders[1]  
    content_2.text = "This is the second slide."  
  
    return prs  
  
def test_extract_text_from_ppt(mock_presentation):  
    with patch('docu_intel.Presentation', return_value=mock_presentation):  
        ppt_file = "mock_presentation.pptx"  
        result = extract_text_from_ppt(ppt_file)  
          
        expected = [  
            {"slide_number": 1, "slide_title": "Slide 1 Title", "text": "This is the first slide."},  
            {"slide_number": 2, "slide_title": "Slide 2 Title", "text": "This is the second slide."}  
        ]  
          
        # Debug print statements  
        print(f"Result: {result}")  
        print(f"Expected: {expected}")  
          
        # Check each item individually  
        for res, exp in zip(result, expected):  
            assert res["slide_number"] == exp["slide_number"]  
            assert res["slide_title"] == exp["slide_title"]  
            assert res["text"].strip() == exp["text"].strip()  
