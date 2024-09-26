import pytest  
from unittest.mock import patch, MagicMock, PropertyMock  
from pptx import Presentation  
import io  
  
# Assuming the detect_image_slides and is_image_of_interest functions are in a module named ppt_image_detector  
# from docu_intel import detect_image_slides, is_image_of_interest  
  
def is_image_of_interest(shape):  
    if hasattr(shape, 'image') and shape.image:  
        # Check the filename extension  
        filename = shape.image.filename.lower()  
        if filename.endswith(('.jpg', '.jpeg', '.png', '.bmp', '.gif')):  
            return filename.split('.')[-1]  
    return None  
  
def detect_image_slides(ppt_bytes):  
    prs = Presentation(io.BytesIO(ppt_bytes))  
    image_slides = {}  
  
    for idx, slide in enumerate(prs.slides, start=1):  
        print(f"Processing slide {idx}")  
        for shape in slide.shapes:  
            print(f"Processing shape {shape}")  
            image_extension = is_image_of_interest(shape)  
            if image_extension:  
                print(f"Found image of interest on slide {idx}: {image_extension}")  
                image_slides[idx] = image_extension  
                break  # We assume only one image of interest per slide for this example  
            else:  
                print("No image of interest found in this shape.")  
      
    print(f"Final detected image slides: {image_slides}")  
    return image_slides  
  
@pytest.fixture  
def input_file_bytes():  
    with open('temp_ppt.pptx', 'rb') as f:  
        return f.read()  
  
@patch('docu_intel.is_image_of_interest')  
def test_detect_image_slides(mock_is_image_of_interest, input_file_bytes):  
    ppt_bytes = input_file_bytes  
  
    # Mock the is_image_of_interest function  
    def side_effect(shape):  
        if hasattr(shape, 'image') and shape.image:  
            if "image1.jpg" in shape.image.filename:  
                return ".jpg"  
            elif "document.pdf" in shape.image.filename:  
                return None  
        return None  
  
    mock_is_image_of_interest.side_effect = side_effect  
  
    with patch('pptx.Presentation', return_value=Presentation(io.BytesIO(ppt_bytes))):  
        result = detect_image_slides(ppt_bytes)  
      
    expected = {  
        7: 'jpeg', 9: 'jpeg'  
    }  
  
    print(result)  
    print(expected)  
  
    assert result == expected  
